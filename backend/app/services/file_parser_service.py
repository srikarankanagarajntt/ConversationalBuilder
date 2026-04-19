"""File parser service — extracts raw text from PDF, DOCX, and PPTX files."""
from __future__ import annotations

import io
import re
from typing import Any, Dict, List

from app.core.exceptions import UnsupportedFileTypeError


class FileParserService:
    async def extract_text(self, raw_bytes: bytes, content_type: str) -> str:
        """
        Dispatch to the appropriate parser based on MIME type.
        Returns the full extracted plain text.
        """
        if "pdf" in content_type:
            return self._parse_pdf(raw_bytes)
        elif "wordprocessingml" in content_type or "msword" in content_type:
            return self._parse_docx(raw_bytes)
        elif "presentationml" in content_type:
            return self._parse_pptx(raw_bytes)
        else:
            raise UnsupportedFileTypeError(content_type)

    def _parse_pdf(self, data: bytes) -> str:
        import PyPDF2

        reader = PyPDF2.PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)

    def _parse_docx(self, data: bytes) -> str:
        from docx import Document

        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)

    def _parse_pptx(self, data: bytes) -> str:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)

    def parse_cv_data(self, text: str) -> Dict[str, Any]:
        """
        Parse CV text and extract structured data.
        """
        lines = text.split('\n')
        lines = [line.strip() for line in lines if line.strip()]

        extracted = {
            "header": self._extract_header(lines),
            "professionalSummary": self._extract_professional_summary(text),
            "technicalSkills": self._extract_technical_skills(text),
            "workExperience": self._extract_work_experience(text),
            "certifications": self._extract_certifications(text)
        }
        return extracted

    def _extract_header(self, lines: List[str]) -> Dict[str, str]:
        header = {"fullName": "", "jobTitle": "", "email": "", "phone": ""}

        # Full name from first line
        if lines:
            header["fullName"] = lines[0]

        # Email regex
        email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
        emails = re.findall(email_pattern, '\n'.join(lines))
        if emails:
            header["email"] = emails[0]

        # Phone regex (simple pattern)
        phone_pattern = r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b'
        phones = re.findall(phone_pattern, '\n'.join(lines))
        if phones:
            header["phone"] = phones[0]

        # Job title - look for common patterns
        title_keywords = ['engineer', 'developer', 'manager', 'analyst', 'architect', 'consultant', 'lead', 'senior', 'junior']
        for line in lines[:10]:  # Check first 10 lines
            lower_line = line.lower()
            if any(keyword in lower_line for keyword in title_keywords):
                header["jobTitle"] = line
                break

        return header

    def _extract_professional_summary(self, text: str) -> List[str]:
        # Look for summary section
        summary_patterns = [
            r'Summary\s*\n(.*?)(?=\n[A-Z][a-z]+|\nTechnical|\nExperience|\nWork|\nSkills|\nEducation|\Z)',
            r'Professional Summary\s*\n(.*?)(?=\n[A-Z][a-z]+|\nTechnical|\nExperience|\nWork|\nSkills|\nEducation|\Z)',
            r'Objective\s*\n(.*?)(?=\n[A-Z][a-z]+|\nTechnical|\nExperience|\nWork|\nSkills|\nEducation|\Z)'
        ]

        for pattern in summary_patterns:
            match = re.search(pattern, text, re.DOTALL | re.IGNORECASE)
            if match:
                summary_text = match.group(1).strip()
                # Split into sentences
                sentences = re.split(r'[.!?]+', summary_text)
                sentences = [s.strip() for s in sentences if s.strip()]
                return sentences[:5]  # Limit to 5 sentences

        return []

    def _extract_technical_skills(self, text: str) -> Dict[str, List[Dict[str, str]]]:
        skills = {"primary": [], "secondary": []}

        # Look for skills section
        skills_pattern = r'Skills?\s*\n(.*?)(?=\n[A-Z][a-z]+|\nExperience|\nWork|\nEducation|\nProjects|\Z)'
        match = re.search(skills_pattern, text, re.DOTALL | re.IGNORECASE)
        if match:
            skills_text = match.group(1).strip()
            # Split by common delimiters
            skill_items = re.split(r'[,;•\-\n]', skills_text)
            skill_items = [s.strip() for s in skill_items if s.strip()]

            # Categorize - first few as primary
            primary_count = min(5, len(skill_items) // 2 + 1)
            for i, skill in enumerate(skill_items):
                skill_dict = {"skill_name": skill, "proficiency": "advanced" if i < primary_count else "intermediate"}
                if i < primary_count:
                    skills["primary"].append(skill_dict)
                else:
                    skills["secondary"].append(skill_dict)

        return skills

    def _extract_work_experience(self, text: str) -> List[Dict[str, str]]:
        experiences = []

        # Look for experience section
        exp_pattern = r'(?:Work )?Experience\s*\n(.*?)(?=\n[A-Z][a-z]+|\nEducation|\nProjects|\nSkills|\Z)'
        match = re.search(exp_pattern, text, re.DOTALL | re.IGNORECASE)
        if match:
            exp_text = match.group(1).strip()
            # Split into entries (rough heuristic)
            entries = re.split(r'\n(?=[A-Z][a-zA-Z\s]+(?:\n|\s*[-–]\s*))', exp_text)
            for entry in entries[:3]:  # Limit to 3
                lines = entry.strip().split('\n')
                if lines:
                    employer = lines[0].strip()
                    position = lines[1].strip() if len(lines) > 1 else ""
                    description = ' '.join(lines[2:]) if len(lines) > 2 else ""
                    experiences.append({
                        "employer": employer,
                        "position": position,
                        "project_description": description
                    })

        return experiences

    def _extract_certifications(self, text: str) -> List[Dict[str, str]]:
        """
        Extract certifications from CV text.
        Looks for sections with 'Certification', 'Certifications', or 'Credentials'.
        """
        certifications = []
        cert_pattern = r'(?:certification|certifications|credentials)[\s\S]*?(?=(?:education|projects|languages|$))'
        cert_sections = re.findall(cert_pattern, text, re.IGNORECASE)
        
        if not cert_sections:
            return certifications
        
        for section in cert_sections:
            lines = section.split('\n')
            for line in lines:
                line = line.strip()
                if not line or line.lower() in ['certification', 'certifications', 'credentials']:
                    continue
                
                # Try to extract name, issuer, and date from line
                # Format could be: "AWS Solutions Architect - Amazon (2023)" or similar
                parts = line.split('-') if '-' in line else [line]
                
                name = parts[0].strip() if len(parts) > 0 else line
                issuer = parts[1].strip() if len(parts) > 1 else ""
                
                # Try to extract date from issuer or from end of name
                date_pattern = r'\((\d{4})\)|\b(20\d{2})\b'
                date_match = re.search(date_pattern, issuer + line)
                date = date_match.group(1) or date_match.group(2) if date_match else ""
                
                if name and name != line:
                    certifications.append({
                        "name": name,
                        "issuer": issuer,
                        "date": date
                    })
                elif name:
                    # If no dash found, try to parse as name - issuer (date)
                    match = re.match(r'(.*?)\s*(?:by|from|-)?\s*([^(]*?)?(?:\(([^)]*)\))?$', line)
                    if match:
                        certifications.append({
                            "name": match.group(1).strip(),
                            "issuer": (match.group(2) or "").strip(),
                            "date": (match.group(3) or "").strip()
                        })
        
        return certifications

