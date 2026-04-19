"""Export service — generates PDF, DOCX, and JSON output files."""
from __future__ import annotations

import json
import os
import re
import subprocess
import tempfile
import threading
import uuid
import asyncio
from datetime import datetime
from typing import Any, Dict, Iterable, List, Tuple

from fastapi import HTTPException, status

from app.models.cv_schema import CvSchema
from app.services.llm_service import LLMService
from app.services.translation_service import TranslationService

# In-memory job store — keyed by job_id
_jobs: Dict[str, Dict[str, Any]] = {}

# Temporary output directory
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "exports")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Master template path
TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "template")
MASTER_TEMPLATE = os.path.join(TEMPLATE_DIR, "NTT_DATA_Master_Templates.docx")


class ExportService:
    def __init__(self):
        self.translation_service = TranslationService()
        self._llm_service = None

    async def create_export_job(self, cv: CvSchema, fmt: str, template_id: str = "ntt-classic", language: str = "en") -> Dict[str, Any]:
        """Generate the file synchronously (POC) and store the job record."""
        job_id = str(uuid.uuid4())
        file_id = str(uuid.uuid4())
        
        # Translate CV if needed
        if language and language != "en":
            try:
                cv = self.translation_service.translate_cv(cv, language)
            except Exception as e:
                # Log but continue with original CV if translation fails
                print(f"Translation failed: {e}")

        try:
            if fmt == "json":
                file_path = self._export_json(cv, file_id)
                media_type = "application/json"
                filename = f"{cv.personalInfo.fullName}_cv.json"
            elif fmt == "docx":
                file_path = self._export_docx(cv, file_id, template_id)
                media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                filename = f"{cv.personalInfo.fullName}_cv.docx"
            elif fmt == "pdf":
                file_path = self._export_pdf(cv, file_id, template_id)
                media_type = "application/pdf"
                filename = f"{cv.personalInfo.fullName}_cv.pdf"
            elif fmt in ("ppt", "pptx"):
                file_path = self._export_ppt(cv, file_id, template_id)
                media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                filename = f"{cv.personalInfo.fullName}_cv.pptx"
            else:
                raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=f"Unknown format: {fmt}")

            job = {
                "jobId": job_id,
                "fileId": file_id,
                "format": fmt,
                "status": "ready",
                "filePath": file_path,
                "mediaType": media_type,
                "filename": filename,
                "downloadUrl": f"/api/download/{file_id}",
            }
            _jobs[job_id] = job
            _jobs[file_id] = job  # also indexed by file_id for downloads
            return job
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Export failed: {str(e)}"
            )

    def get_job(self, job_id: str) -> Dict[str, Any]:
        job = _jobs.get(job_id)
        if not job:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=f"Job '{job_id}' not found.")
        return job

    def get_file_path(self, file_id: str) -> Tuple[str, str, str]:
        job = _jobs.get(file_id)
        if not job:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found.")
        return job["filePath"], job["mediaType"], job["filename"]

    # ── Format generators ─────────────────────────────────────────────────────

    def _export_json(self, cv: CvSchema, file_id: str) -> str:
        """Export CV as JSON."""
        path = os.path.join(OUTPUT_DIR, f"{file_id}.json")
        with open(path, "w", encoding="utf-8") as f:
            json.dump(cv.model_dump(), f, indent=2)
        return path

    def _export_docx(self, cv: CvSchema, file_id: str, template_id: str = "ntt-classic") -> str:
        """Export CV as DOCX using template rendering."""
        try:
            # Render from template
            rendered_docx_path = self._render_docx_from_template(cv, template_id)

            # Verify the rendered file exists and has content
            if not os.path.exists(rendered_docx_path):
                raise FileNotFoundError(f"Rendered DOCX file not created: {rendered_docx_path}")

            file_size = os.path.getsize(rendered_docx_path)
            if file_size == 0:
                raise ValueError("Rendered DOCX file is empty")

            # Copy to output directory
            output_path = os.path.join(OUTPUT_DIR, f"{file_id}.docx")
            with open(rendered_docx_path, "rb") as src:
                content = src.read()
                if not content:
                    raise ValueError("Failed to read rendered DOCX file")
                with open(output_path, "wb") as dst:
                    dst.write(content)

            # Verify the output file
            if not os.path.exists(output_path):
                raise FileNotFoundError(f"Output DOCX file not created: {output_path}")

            output_size = os.path.getsize(output_path)
            if output_size == 0:
                raise ValueError("Output DOCX file is empty")

            # Cleanup temporary file
            if os.path.exists(rendered_docx_path):
                try:
                    os.remove(rendered_docx_path)
                except Exception as cleanup_error:
                    print(f"Warning: Could not clean up temporary file {rendered_docx_path}: {cleanup_error}")

            return output_path
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"DOCX rendering failed: {str(e)}"
            )

    def _export_pdf(self, cv: CvSchema, file_id: str, template_id: str = "ntt-classic") -> str:
        """Export CV as PDF (via DOCX rendering + conversion)."""
        try:
            # First render DOCX
            docx_path = self._render_docx_from_template(cv, template_id)

            # Then convert to PDF
            pdf_path = os.path.join(OUTPUT_DIR, f"{file_id}.pdf")
            self._convert_docx_to_pdf(docx_path, pdf_path)

            # Cleanup temporary DOCX
            if os.path.exists(docx_path):
                os.remove(docx_path)

            return pdf_path
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"PDF conversion failed: {str(e)}"
            )

    def _render_docx_from_template(self, cv: CvSchema, template_id: str) -> str:
        """Render DOCX from master template with Jinja2 placeholders."""
        try:
            from docxtpl import DocxTemplate
        except ImportError:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="docxtpl library not installed. Run: pip install docxtpl"
            )

        # Check if master template exists
        if not os.path.exists(MASTER_TEMPLATE):
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Master template not found: {MASTER_TEMPLATE}"
            )

        try:
            # Load template
            doc = DocxTemplate(MASTER_TEMPLATE)

            # Prepare context for Jinja2 rendering
            context = self._prepare_template_context(cv)
            
            # Ensure all strings in context are properly encoded for DOCX
            context = self._sanitize_context_for_docx(context)

            # Render template
            doc.render(context)

            # Save to temporary file
            temp_file = os.path.join(tempfile.gettempdir(), f"rendered_{uuid.uuid4()}.docx")
            doc.save(temp_file)

            return temp_file
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"DOCX rendering failed: {str(e)}"
            )

    def _sanitize_context_for_docx(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Sanitize context values to ensure proper DOCX encoding."""
        sanitized = {}
        for key, value in context.items():
            if isinstance(value, str):
                # Ensure UTF-8 encoding and remove any null characters
                try:
                    # Clean up any problematic characters
                    cleaned = value.replace('\x00', '').replace('\r\n', '\n')
                    # Ensure it can be encoded to UTF-8
                    cleaned.encode('utf-8')
                    sanitized[key] = cleaned
                except (UnicodeEncodeError, AttributeError):
                    # Fallback to removing problematic chars
                    sanitized[key] = value.encode('utf-8', errors='ignore').decode('utf-8')
            elif isinstance(value, list):
                # Recursively sanitize list items
                sanitized[key] = [
                    self._sanitize_context_for_docx(item) if isinstance(item, dict)
                    else (item.encode('utf-8', errors='ignore').decode('utf-8') if isinstance(item, str) else item)
                    for item in value
                ]
            elif isinstance(value, dict):
                # Recursively sanitize nested dictionaries
                sanitized[key] = self._sanitize_context_for_docx(value)
            else:
                sanitized[key] = value
        return sanitized

    def _prepare_template_context(self, cv: CvSchema) -> Dict[str, Any]:
        """Prepare CV data for Jinja2 template rendering."""
        return {
            # Personal Info
            "fullName": cv.personalInfo.fullName or "",
            "jobTitle": cv.header.get("jobTitle") or "",
            "email": cv.personalInfo.email or "",
            "phone": cv.personalInfo.phone or "",
            "location": cv.personalInfo.location or "",
            "summary": cv.personalInfo.summary or "",

            # Skills
            "skills": cv.skills or [],

            # Experience
            "experience": [
                {
                    "company": exp.company,
                    "title": exp.title,
                    "role": exp.role,
                    "startDate": exp.startDate,
                    "endDate": exp.endDate,
                    "location": exp.location,
                    "projectName": exp.projectName,
                    "projectInformation": exp.projectInformation,
                    "clients": exp.clients,
                    "technology": exp.technology or [],
                    "description": exp.description,
                    "achievements": exp.achievements or []
                }
                for exp in cv.experience
            ] if cv.experience else [],

            # Education
            "education": [
                {
                    "institution": edu.institution,
                    "degree": edu.degree,
                    "field": edu.field,
                    "startDate": edu.startDate,
                    "endDate": edu.endDate
                }
                for edu in cv.education
            ] if cv.education else [],

            # Projects
            "projects": [
                {
                    "name": proj.name,
                    "description": proj.description,
                    "technologies": proj.technologies or [],
                    "url": proj.url
                }
                for proj in cv.projects
            ] if cv.projects else [],

            # Certifications
            "certifications": [
                {
                    "name": cert.name,
                    "issuer": cert.issuer,
                    "date": cert.date
                }
                for cert in cv.certifications
            ] if cv.certifications else [],

            # Languages
            "languages": cv.languages or [],

            # Technical Skills with proficiency
            "technicalSkills": {
                "primary": cv.technicalSkills.get("primary", []) if cv.technicalSkills else [],
                "secondary": cv.technicalSkills.get("secondary", []) if cv.technicalSkills else []
            },

            # Professional Summary
            "professionalSummary": cv.professionalSummary or [],

            # Header info
            "header": cv.header or {}
        }

    def _convert_docx_to_pdf(self, docx_path: str, pdf_path: str) -> None:
        """Convert DOCX to PDF using available tools."""
        if not os.path.exists(docx_path):
            raise FileNotFoundError(f"DOCX file not found: {docx_path}")

        # Try docx2pdf first (Windows/macOS with Word installed)
        if self._try_docx2pdf_conversion(docx_path, pdf_path):
            return

        # Fallback to LibreOffice
        if self._try_libreoffice_conversion(docx_path, pdf_path):
            return

        # All methods failed
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="PDF conversion failed: docx2pdf or LibreOffice required"
        )

    def _try_docx2pdf_conversion(self, docx_path: str, pdf_path: str) -> bool:
        """Try to convert using docx2pdf library."""
        try:
            from docx2pdf import convert
            convert(docx_path, pdf_path)
            return os.path.exists(pdf_path)
        except ImportError:
            return False
        except Exception as e:
            print(f"docx2pdf conversion failed: {e}")
            return False

    def _try_libreoffice_conversion(self, docx_path: str, pdf_path: str) -> bool:
        """Try to convert using LibreOffice headless conversion."""
        try:
            output_dir = os.path.dirname(pdf_path)

            # Build command for LibreOffice
            cmd = [
                "soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                docx_path
            ]

            # Run conversion
            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=60
            )

            # Check if conversion succeeded
            if result.returncode == 0 and os.path.exists(pdf_path):
                return True

            return False
        except FileNotFoundError:
            # LibreOffice not installed
            return False
        except subprocess.TimeoutExpired:
            print("LibreOffice conversion timed out")
            return False
        except Exception as e:
            print(f"LibreOffice conversion error: {e}")
            return False

    def _export_ppt(self, cv, file_id, template_id):
        """Export CV as PPTX using the latest one-slide template."""
        try:
            from pptx import Presentation
        except ImportError:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="python-pptx library not installed. Run: pip install python-pptx",
            )

        try:
            template_path = self._resolve_ppt_template_path(template_id)
            presentation = Presentation(template_path)

            normalized = self._normalize_cv_for_ppt(cv)
            generated_month_year = datetime.now().strftime("%B %Y")

            summary_text = self._compact_text_for_slide(
                normalized["summary"],
                max_chars=320,
                text_kind="professional summary",
            )
            primary_skills_text = self._compact_skills_list_for_slide(
                normalized["primary_skills"],
                max_chars=140,
                text_kind="primary skills",
            )
            secondary_skills_text = self._compact_skills_list_for_slide(
                normalized["secondary_skills"],
                max_chars=140,
                text_kind="secondary skills",
            )
            certifications_text = ", ".join(normalized["certifications"])
            languages_text = normalized["languages"]
            experience_text = self._build_experience_block(normalized["experiences"])
            header_text = self._build_header_text(normalized["job_title"], normalized["full_name"])
            profile_text = self._build_profile_text(
                normalized["email"],
                normalized["phone"],
                normalized["location"],
                summary_text,
                languages_text,
            )
            skills_block_text = self._build_skills_block_text(
                primary_skills_text,
                secondary_skills_text,
                certifications_text,
            )

            replacements = {
                "{{generatedMonthYear}}": generated_month_year,
                "{{cvDraft.header.jobTitle}}": normalized["job_title"],
                "{{cvDraft.header.fullName}}": normalized["full_name"],
                "{{cvDraft.personalInfo.summary}}": summary_text,
                "{{summary}}": summary_text,
                "{{email}}": normalized["email"],
                "{{phone}}": normalized["phone"],
                "{{location}}": normalized["location"],
                "{{languages}}": languages_text,
                "{{cvDraft.technicalSkills.primary[*].skillName}}": primary_skills_text,
                "{{cvDraft.technicalSkills.secondary[*].skillName}}": secondary_skills_text,
                "{{cvDraft.skills}}": secondary_skills_text,
                "{{cvDraft.certifications}}": certifications_text,
            }

            for idx in range(1, 5):
                exp = normalized["experiences"][idx - 1] if idx <= len(normalized["experiences"]) else {}
                company = exp.get("company", "")
                title = exp.get("title", "")
                description = exp.get("description", "")
                date_range = self._build_date_range(exp.get("start_date", ""), exp.get("end_date", ""))

                replacements[f"{{{{company{idx}}}}}"] = company
                replacements[f"{{{{company{idx}_ttile}}}}"] = title
                replacements[f"{{{{company{idx}_project_description}}}}"] = description
                replacements[f"{{{{company{idx}_start_Date-company{idx}_end_date}}}}"] = date_range

                replacements[f"{{{{workExperience{idx}.company}}}}"] = company
                replacements[f"{{{{workExperience{idx}.title}}}}"] = title
                replacements[f"{{{{workExperience{idx}.description}}}}"] = description
                replacements[f"{{{{workExperience{idx}.dateRange}}}}"] = date_range

            shape_fallbacks = {
                "generated_date_box": generated_month_year,
                "header_box": header_text,
                "profile_box": profile_text,
                "skills_box": skills_block_text,
                "experience_box": experience_text,
            }

            self._replace_ppt_placeholders(
                presentation,
                replacements,
                experience_text,
                shape_fallbacks,
                normalized["experiences"],
                {
                    "contact": " | ".join([value for value in [normalized["email"], normalized["phone"], normalized["location"]] if value]),
                    "summary": summary_text,
                    "languages": languages_text,
                },
                {
                    "primary_skills": primary_skills_text,
                    "secondary_skills": secondary_skills_text,
                    "certifications": certifications_text,
                },
            )

            output_path = os.path.join(OUTPUT_DIR, f"{file_id}.pptx")
            presentation.save(output_path)
            return output_path
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"PPT export failed: {str(e)}",
            )

    def _resolve_ppt_template_path(self, template_id: str) -> str:
        """Resolve PPT template file from template_id with fallback to latest master."""
        default_template_name = "Resume_master_with_placeholders.pptx"
        template_map = {
            "java-architect-pptx": default_template_name,
            "ntt-classic": default_template_name,
            "ntt-modern": default_template_name,
            "minimal": default_template_name,
        }

        candidates: List[str] = []
        if template_id:
            if os.path.isabs(template_id):
                candidates.append(template_id)
            if template_id.lower().endswith(".pptx"):
                candidates.append(os.path.join(TEMPLATE_DIR, template_id))
            mapped = template_map.get(template_id)
            if mapped:
                candidates.append(os.path.join(TEMPLATE_DIR, mapped))
            candidates.append(os.path.join(TEMPLATE_DIR, f"{template_id}.pptx"))

        candidates.append(os.path.join(TEMPLATE_DIR, default_template_name))

        for candidate in candidates:
            if candidate and os.path.exists(candidate):
                return candidate

        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"PPT template not found for template_id '{template_id}'",
        )

    def _normalize_cv_for_ppt(self, cv: Any) -> Dict[str, Any]:
        """Normalize incoming CV payload into stable fields for PPT rendering."""
        if hasattr(cv, "model_dump"):
            source = cv.model_dump(by_alias=True)
        elif isinstance(cv, dict):
            source = cv
        else:
            source = {}

        cv_draft = source.get("cvDraft", source)
        header = cv_draft.get("header", {}) or {}
        header_contact = header.get("contact", {}) or {}
        personal_info = cv_draft.get("personalInfo", {}) or {}

        full_name = self._cleanup_text(
            header.get("fullName") or header.get("name") or personal_info.get("fullName") or personal_info.get("name") or "",
            preserve_newlines=False,
        )
        job_title = self._cleanup_text(
            header.get("jobTitle") or header.get("title") or personal_info.get("jobTitle") or personal_info.get("title") or "",
            preserve_newlines=False,
        )
        email = self._cleanup_text(
            header.get("email") or header_contact.get("email") or header_contact.get("emailId") or personal_info.get("email") or "",
            preserve_newlines=False,
        )
        phone = self._cleanup_text(
            header.get("phone") or header_contact.get("phone") or personal_info.get("phone") or "",
            preserve_newlines=False,
        )
        location = self._cleanup_text(
            personal_info.get("location") or header_contact.get("location") or "",
            preserve_newlines=False,
        )

        summary = personal_info.get("summary") or ""
        if not summary:
            summary = " ".join([str(item).strip() for item in (cv_draft.get("professionalSummary") or []) if str(item).strip()])
        summary = self._cleanup_text(summary)

        technical_skills = cv_draft.get("technicalSkills", {}) or {}
        primary_skills = self._extract_skill_names(technical_skills.get("primary", []))
        secondary_skills = self._extract_skill_names(technical_skills.get("secondary", []))
        fallback_skills = self._extract_skill_names(cv_draft.get("skills", []))

        merged_skills: List[str] = []
        seen = set()
        for skill in primary_skills + secondary_skills + fallback_skills:
            key = skill.casefold()
            if not key or key in seen:
                continue
            seen.add(key)
            merged_skills.append(skill)

        normalized_primary: List[str] = []
        seen_primary = set()
        for skill in primary_skills:
            key = skill.casefold()
            if not key or key in seen_primary:
                continue
            seen_primary.add(key)
            normalized_primary.append(skill)

        normalized_secondary: List[str] = []
        seen_secondary = set()
        for skill in secondary_skills + fallback_skills:
            key = skill.casefold()
            if not key or key in seen_secondary or key in seen_primary:
                continue
            seen_secondary.add(key)
            normalized_secondary.append(skill)

        if not normalized_primary and merged_skills:
            normalized_primary = merged_skills[:8]
            seen_primary = {skill.casefold() for skill in normalized_primary}

        if not normalized_secondary:
            normalized_secondary = [skill for skill in merged_skills if skill.casefold() not in seen_primary][:8]

        certifications = self._extract_certification_names(cv_draft.get("certifications", []))
        languages = self._serialize_languages(
            cv_draft.get("languages")
            or cv_draft.get("personalDetails", {}).get("languages")
            or []
        )

        project_descriptions = self._extract_project_descriptions(cv_draft.get("projects") or [])
        experiences = self._extract_latest_work_experiences(cv_draft, project_descriptions)

        return {
            "full_name": full_name,
            "job_title": job_title,
            "email": email,
            "phone": phone,
            "location": location,
            "summary": summary,
            "skills": merged_skills,
            "primary_skills": normalized_primary,
            "secondary_skills": normalized_secondary,
            "certifications": certifications,
            "languages": languages,
            "experiences": experiences,
        }

    def _extract_skill_names(self, skills: Any) -> List[str]:
        values: List[str] = []
        if not isinstance(skills, list):
            return values

        for item in skills:
            if isinstance(item, str):
                skill = self._cleanup_text(item, preserve_newlines=False)
            elif isinstance(item, dict):
                skill = self._cleanup_text(
                    item.get("skillName")
                    or item.get("skill_name")
                    or item.get("name")
                    or "",
                    preserve_newlines=False,
                )
            else:
                skill = ""
            if skill:
                values.append(skill)
        return values

    def _extract_certification_names(self, certifications: Any) -> List[str]:
        names: List[str] = []
        if not isinstance(certifications, list):
            return names
        for cert in certifications:
            if isinstance(cert, str):
                name = self._cleanup_text(cert, preserve_newlines=False)
            elif isinstance(cert, dict):
                name = self._cleanup_text(cert.get("name") or cert.get("title") or "", preserve_newlines=False)
            else:
                name = ""
            if name:
                names.append(name)
        return names

    def _extract_project_descriptions(self, projects: Any) -> List[str]:
        descriptions: List[str] = []
        if not isinstance(projects, list):
            return descriptions

        for project in projects:
            if isinstance(project, dict):
                description = self._cleanup_text(
                    project.get("description") or project.get("projectDescription") or "",
                    preserve_newlines=False,
                )
            else:
                description = ""
            if description:
                descriptions.append(description)
        return descriptions

    def _serialize_languages(self, languages: Any) -> str:
        if not isinstance(languages, list):
            return ""
        values = []
        for language in languages:
            if isinstance(language, str):
                lang = self._cleanup_text(language, preserve_newlines=False)
            elif isinstance(language, dict):
                lang = self._cleanup_text(
                    language.get("name") or language.get("language") or "",
                    preserve_newlines=False,
                )
                proficiency = self._cleanup_text(language.get("proficiency") or "", preserve_newlines=False)
                if lang and proficiency:
                    lang = f"{lang} ({proficiency})"
            else:
                lang = ""
            if lang:
                values.append(lang)
        return ", ".join(values)

    def _extract_latest_work_experiences(self, cv_draft: Dict[str, Any], project_descriptions: List[str] | None = None) -> List[Dict[str, str]]:
        """Prefer workExperience, fallback to experience, and keep latest 4."""
        raw_experience = cv_draft.get("workExperience") or cv_draft.get("experience") or []
        if not isinstance(raw_experience, list):
            return []

        project_descriptions = project_descriptions or []
        normalized = [
            self._normalize_experience_entry(
                entry,
                fallback_project_description=project_descriptions[idx] if idx < len(project_descriptions) else "",
            )
            for idx, entry in enumerate(raw_experience)
        ]
        normalized = [exp for exp in normalized if exp["company"] or exp["title"] or exp["description"]]

        if not normalized:
            return []

        if any(self._date_sort_score(exp.get("end_date", ""), exp.get("start_date", "")) > -1 for exp in normalized):
            normalized.sort(
                key=lambda exp: self._date_sort_score(exp.get("end_date", ""), exp.get("start_date", "")),
                reverse=True,
            )

        return normalized[:4]

    def _normalize_experience_entry(self, entry: Any, fallback_project_description: str = "") -> Dict[str, str]:
        if hasattr(entry, "model_dump"):
            entry = entry.model_dump(by_alias=True)
        elif not isinstance(entry, dict) and hasattr(entry, "__dict__"):
            entry = vars(entry)
        if not isinstance(entry, dict):
            entry = {}

        achievements = entry.get("achievements") or []
        if isinstance(achievements, list):
            achievements_text = " ".join([str(item).strip() for item in achievements if str(item).strip()])
        else:
            achievements_text = str(achievements).strip()
        responsibilities = entry.get("responsibilities") or []
        if isinstance(responsibilities, list):
            responsibilities_text = " ".join([str(item).strip() for item in responsibilities if str(item).strip()])
        else:
            responsibilities_text = str(responsibilities).strip()

        description = (
            entry.get("projectDescription")
            or entry.get("project_description")
            or entry.get("description")
            or achievements_text
            or responsibilities_text
            or fallback_project_description
            or ""
        )

        normalized = {
            "company": self._cleanup_text(
                entry.get("employer")
                or entry.get("company")
                or entry.get("companyName")
                or entry.get("organization")
                or entry.get("client")
                or "",
                preserve_newlines=False,
            ),
            "title": self._cleanup_text(
                entry.get("position")
                or entry.get("title")
                or entry.get("jobTitle")
                or entry.get("role")
                or entry.get("designation")
                or "",
                preserve_newlines=False,
            ),
            "description": self._compact_text_for_slide(description, max_chars=180, text_kind="experience description"),
            "start_date": self._cleanup_text(
                entry.get("startDate") or entry.get("start_date") or entry.get("start") or entry.get("from") or "",
                preserve_newlines=False,
            ),
            "end_date": self._cleanup_text(
                entry.get("endDate") or entry.get("end_date") or entry.get("end") or entry.get("to") or "",
                preserve_newlines=False,
            ),
        }
        return normalized

    def _build_date_range(self, start_date: str, end_date: str) -> str:
        start = self._cleanup_text(start_date, preserve_newlines=False)
        end = self._cleanup_text(end_date, preserve_newlines=False)
        if start and end:
            return f"{start} - {end}"
        return start or end

    def _date_sort_score(self, end_date: str, start_date: str) -> int:
        def _extract_score(value: str) -> int:
            text = (value or "").strip().lower()
            if not text:
                return -1
            if any(token in text for token in ("present", "current", "now", "ongoing")):
                return 999999
            years = re.findall(r"(?:19|20)\d{2}", text)
            if not years:
                return -1
            return int(years[-1])

        end_score = _extract_score(end_date)
        if end_score != -1:
            return end_score
        return _extract_score(start_date)

    def _compact_skills_for_slide(self, skills: List[str]) -> str:
        if not skills:
            return ""

        compact = ", ".join(skills[:20])
        compact = self._cleanup_text(compact, preserve_newlines=False)
        if len(compact) > 220:
            optimized = self._optimize_text_with_llm(compact, 220, "skills section")
            if optimized:
                compact = self._cleanup_text(optimized, preserve_newlines=False)
        if len(compact) > 220:
            compact = self._truncate_with_ellipsis(compact, 220)
        return compact

    def _compact_primary_skills_for_profile(self, skills: List[str]) -> str:
        if not skills:
            return ""
        compact = ", ".join(skills[:10])
        compact = self._cleanup_text(compact, preserve_newlines=False)
        if len(compact) > 140:
            compact = self._truncate_with_ellipsis(compact, 140)
        return compact

    def _compact_skills_list_for_slide(self, skills: List[str], max_chars: int, text_kind: str) -> str:
        if not skills:
            return ""

        compact = ", ".join(skills[:12])
        compact = self._cleanup_text(compact, preserve_newlines=False)
        if len(compact) > max_chars:
            optimized = self._optimize_text_with_llm(compact, max_chars, text_kind)
            if optimized:
                compact = self._cleanup_text(optimized, preserve_newlines=False)
        if len(compact) > max_chars:
            compact = self._truncate_with_ellipsis(compact, max_chars)
        return compact

    def _compact_text_for_slide(self, text: str, max_chars: int, text_kind: str) -> str:
        cleaned = self._cleanup_text(text)
        if len(cleaned) <= max_chars:
            return cleaned

        lightly_clipped = self._truncate_with_ellipsis(cleaned, max_chars + 40)
        if len(lightly_clipped) <= max_chars:
            return lightly_clipped

        optimized = self._optimize_text_with_llm(lightly_clipped, max_chars, text_kind)
        if optimized:
            optimized = self._cleanup_text(optimized)
            if len(optimized) <= max_chars:
                return optimized
            return self._truncate_with_ellipsis(optimized, max_chars)

        return self._truncate_with_ellipsis(lightly_clipped, max_chars)

    def _cleanup_text(self, value: Any, preserve_newlines: bool = True) -> str:
        text = str(value or "")
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\s*([,;:])\s*", r"\1 ", text)
        text = re.sub(r"\s{2,}", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = text.replace("Responsible for ", "")
        text = text.replace("responsible for ", "")
        text = text.strip()
        if not preserve_newlines:
            text = re.sub(r"\s*\n\s*", " ", text)
            text = re.sub(r"\s{2,}", " ", text).strip()
        return text

    def _truncate_with_ellipsis(self, text: str, max_chars: int) -> str:
        if len(text) <= max_chars:
            return text
        if max_chars <= 3:
            return text[:max_chars]
        return text[: max_chars - 3].rstrip() + "..."

    def _build_experience_block(self, experiences: List[Dict[str, str]]) -> str:
        if not experiences:
            return ""

        lines: List[str] = []
        for exp in experiences[:4]:
            header_parts = [part for part in [exp.get("company", ""), exp.get("title", ""), self._build_date_range(exp.get("start_date", ""), exp.get("end_date", ""))] if part]
            if header_parts:
                lines.append(" | ".join(header_parts))
            if exp.get("description"):
                lines.append(f"     {exp['description']}")

        return "\n".join(lines).strip()

    def _build_header_text(self, job_title: str, full_name: str) -> str:
        lines = [value for value in [self._cleanup_text(job_title, preserve_newlines=False), self._cleanup_text(full_name, preserve_newlines=False)] if value]
        return "\n".join(lines)

    def _build_profile_text(self, email: str, phone: str, location: str, summary: str, languages: str) -> str:
        contact = " | ".join([value for value in [email, phone, location] if value])
        lines: List[str] = []
        if contact:
            lines.append(contact)
        if summary:
            lines.append(summary)
        if languages:
            lines.append(f"LANGUAGES: {languages}")
        return "\n".join(lines).strip()

    def _build_skills_block_text(self, primary_skills_text: str, secondary_skills_text: str, certifications_text: str) -> str:
        lines = ["[PRIMARY]"]
        lines.append(primary_skills_text or "")
        lines.append("[SECONDARY]")
        lines.append(secondary_skills_text or "")
        lines.append("[CERTIFICATIONS]")
        lines.append(certifications_text or "")
        return "\n".join(lines).strip()

    def _replace_ppt_placeholders(
        self,
        presentation,
        replacements: Dict[str, str],
        experience_block: str,
        shape_fallbacks: Dict[str, str],
        experiences: List[Dict[str, str]],
        profile_data: Dict[str, str],
        skills_data: Dict[str, str],
    ) -> None:
        for slide in presentation.slides:
            for shape in self._iter_shapes_recursive(slide.shapes):
                if not getattr(shape, "has_text_frame", False):
                    continue

                original_text = shape.text_frame.text or ""
                if not original_text:
                    continue

                shape_name = (getattr(shape, "name", "") or "").strip()

                if shape_name == "experience_box":
                    self._set_experience_shape_with_bullets(shape, experiences)
                    self._apply_text_style(shape)
                    continue
                elif shape_name == "profile_box":
                    self._set_profile_shape(shape, profile_data)
                    continue
                elif shape_name == "skills_box":
                    self._set_skills_shape(shape, skills_data)
                    continue
                elif shape_name in shape_fallbacks:
                    updated_text = self._normalize_slide_text(shape_fallbacks.get(shape_name, ""))
                elif self._contains_experience_template_block(original_text):
                    updated_text = experience_block
                else:
                    updated_text = original_text
                    for placeholder, value in replacements.items():
                        updated_text = updated_text.replace(placeholder, value or "")
                    updated_text = self._strip_unresolved_placeholders(updated_text)
                    updated_text = self._normalize_slide_text(updated_text)

                if updated_text != original_text:
                    self._set_shape_text(shape, updated_text)
                    if shape_name == "header_box":
                        self._apply_text_style(shape, font_size_pt=16, bold=True)
                    else:
                        self._apply_text_style(shape)
                elif shape_name in shape_fallbacks:
                    if shape_name == "header_box":
                        self._apply_text_style(shape, font_size_pt=16, bold=True)
                    else:
                        self._apply_text_style(shape)

    def _set_experience_shape_with_bullets(self, shape: Any, experiences: List[Dict[str, str]]) -> None:
        """Render professional experience as bullet points."""
        text_frame = getattr(shape, "text_frame", None)
        if not text_frame:
            return

        text_frame.clear()
        if not experiences:
            text_frame.text = ""
            return

        first = True
        for exp in experiences[:4]:
            company = self._cleanup_text(exp.get("company", ""), preserve_newlines=False)
            title = self._cleanup_text(exp.get("title", ""), preserve_newlines=False)
            date_range = self._build_date_range(exp.get("start_date", ""), exp.get("end_date", ""))
            description = self._cleanup_text(exp.get("description", ""), preserve_newlines=False)

            header_parts = [part for part in [company, title, date_range] if part]
            if not header_parts and not description:
                continue

            header_line = " | ".join(header_parts)
            if first:
                paragraph = text_frame.paragraphs[0]
                first = False
            else:
                paragraph = text_frame.add_paragraph()
            paragraph.text = f"• {header_line}" if header_line else "•"
            for run in paragraph.runs:
                if run and run.font:
                    run.font.bold = True

            if description:
                desc_paragraph = text_frame.add_paragraph()
                desc_paragraph.text = f"  {description}"
                for run in desc_paragraph.runs:
                    if run and run.font:
                        run.font.bold = False

    def _set_profile_shape(self, shape: Any, profile_data: Dict[str, str]) -> None:
        """Render profile block with bold contact details."""
        text_frame = getattr(shape, "text_frame", None)
        if not text_frame:
            return

        contact = self._cleanup_text(profile_data.get("contact", ""), preserve_newlines=False)
        summary = self._cleanup_text(profile_data.get("summary", ""))
        languages = self._cleanup_text(profile_data.get("languages", ""), preserve_newlines=False)

        text_frame.clear()
        first = True

        def _next_paragraph():
            nonlocal first
            if first:
                first = False
                return text_frame.paragraphs[0]
            return text_frame.add_paragraph()

        if contact:
            p = _next_paragraph()
            p.text = contact
            self._apply_paragraph_style(p, font_size_pt=9, bold=True)

        if summary:
            p = _next_paragraph()
            p.text = summary
            self._apply_paragraph_style(p, font_size_pt=9, bold=False)

        if languages:
            p = _next_paragraph()
            p.text = f"LANGUAGES: {languages}"
            self._apply_paragraph_style(p, font_size_pt=9, bold=False)

    def _set_skills_shape(self, shape: Any, skills_data: Dict[str, str]) -> None:
        """Render skills block with bold section labels."""
        text_frame = getattr(shape, "text_frame", None)
        if not text_frame:
            return

        primary_skills = self._cleanup_text(skills_data.get("primary_skills", ""), preserve_newlines=False)
        secondary_skills = self._cleanup_text(skills_data.get("secondary_skills", ""), preserve_newlines=False)
        certifications = self._cleanup_text(skills_data.get("certifications", ""), preserve_newlines=False)

        text_frame.clear()
        first = True

        def _next_paragraph():
            nonlocal first
            if first:
                first = False
                return text_frame.paragraphs[0]
            return text_frame.add_paragraph()

        p = _next_paragraph()
        p.text = "PRIMARY:"
        self._apply_paragraph_style(p, font_size_pt=9, bold=True)
        if primary_skills:
            p = _next_paragraph()
            p.text = primary_skills
            self._apply_paragraph_style(p, font_size_pt=9, bold=False)

        p = _next_paragraph()
        p.text = "SECONDARY:"
        self._apply_paragraph_style(p, font_size_pt=9, bold=True)
        if secondary_skills:
            p = _next_paragraph()
            p.text = secondary_skills
            self._apply_paragraph_style(p, font_size_pt=9, bold=False)

        p = _next_paragraph()
        p.text = "CERTIFICATIONS:"
        self._apply_paragraph_style(p, font_size_pt=9, bold=True)
        if certifications:
            p = _next_paragraph()
            p.text = certifications
            self._apply_paragraph_style(p, font_size_pt=9, bold=False)

    def _apply_paragraph_style(self, paragraph: Any, font_size_pt: int = 9, bold: bool = False) -> None:
        try:
            from pptx.dml.color import RGBColor
            from pptx.util import Pt
        except Exception:
            return

        if paragraph and paragraph.font:
            paragraph.font.name = "Calibri"
            paragraph.font.size = Pt(font_size_pt)
            paragraph.font.bold = bool(bold)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)

        for run in paragraph.runs:
            if run and run.font:
                run.font.name = "Calibri"
                run.font.size = Pt(font_size_pt)
                run.font.bold = bool(bold)
                run.font.color.rgb = RGBColor(0, 0, 0)

    def _apply_text_style(self, shape: Any, font_size_pt: int = 9, bold: Any = None) -> None:
        """Apply Calibri 9pt black style for generated text."""
        try:
            from pptx.dml.color import RGBColor
            from pptx.util import Pt
        except Exception:
            return

        text_frame = getattr(shape, "text_frame", None)
        if not text_frame:
            return

        for paragraph in text_frame.paragraphs:
            if paragraph and paragraph.font:
                paragraph.font.name = "Calibri"
                paragraph.font.size = Pt(font_size_pt)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                if bold is not None:
                    paragraph.font.bold = bool(bold)
            for run in paragraph.runs:
                if run and run.font:
                    run.font.name = "Calibri"
                    run.font.size = Pt(font_size_pt)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    if bold is not None:
                        run.font.bold = bool(bold)

    def _contains_experience_template_block(self, text: str) -> bool:
        return (
            "{{company1}}" in text and "{{company2}}" in text
        ) or (
            "{{workExperience1.company}}" in text and "{{workExperience2.company}}" in text
        )

    def _iter_shapes_recursive(self, shapes) -> Iterable[Any]:
        for shape in shapes:
            yield shape
            if hasattr(shape, "shapes"):
                yield from self._iter_shapes_recursive(shape.shapes)

    def _set_shape_text(self, shape: Any, text: str) -> None:
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.text = text or ""

    def _strip_unresolved_placeholders(self, text: str) -> str:
        return re.sub(r"\{\{[^{}]+\}\}", "", text)

    def _normalize_slide_text(self, text: str) -> str:
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.split("\n")]
        cleaned_lines: List[str] = []
        for line in lines:
            if not line:
                if cleaned_lines and cleaned_lines[-1] != "":
                    cleaned_lines.append("")
                continue
            if line in ("|", "-", "| |", "| | |"):
                continue
            cleaned_lines.append(line)
        while cleaned_lines and cleaned_lines[-1] == "":
            cleaned_lines.pop()
        return "\n".join(cleaned_lines)

    def _optimize_text_with_llm(self, text: str, max_chars: int, text_kind: str) -> str:
        if not text or len(text) <= max_chars:
            return text

        if self._llm_service is None:
            try:
                self._llm_service = LLMService()
            except Exception as e:
                print(f"LLM unavailable for PPT optimization: {e}")
                return ""

        prompt = (
            f"Rewrite the following {text_kind} for a one-slide resume PowerPoint.\n"
            "Rules:\n"
            "- Do not invent facts.\n"
            "- Do not add fake metrics.\n"
            "- Do not add new companies, roles, or technologies.\n"
            "- Preserve meaning and strongest keywords.\n"
            "- Use concise professional resume language.\n"
            f"- Keep response under {max_chars} characters.\n"
            "- Return plain text only.\n\n"
            f"Text:\n{text}"
        )

        messages = [
            {
                "role": "system",
                "content": "You optimize resume text for tight slide layouts. Return plain text only.",
            },
            {"role": "user", "content": prompt},
        ]

        try:
            response = self._run_coroutine_sync(self._llm_service.chat(messages))
            return self._cleanup_text(response)
        except Exception as e:
            print(f"LLM optimization failed ({text_kind}): {e}")
            return ""

    def _run_coroutine_sync(self, coroutine):
        try:
            asyncio.get_running_loop()
        except RuntimeError:
            return asyncio.run(coroutine)

        result: Dict[str, Any] = {}
        error: Dict[str, Exception] = {}

        def _runner():
            try:
                result["value"] = asyncio.run(coroutine)
            except Exception as exc:
                error["value"] = exc

        thread = threading.Thread(target=_runner, daemon=True)
        thread.start()
        thread.join()

        if "value" in error:
            raise error["value"]
        return result.get("value", "")

